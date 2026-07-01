"""Deck EXECUTIVO como .pptx TOTALMENTE EDITÁVEL (consulting clean). Kicker, títulos,
réguas, chips de status, callouts, pipeline (cards numerados), dashboard (3 colunas) e
tabelas são objetos nativos; só os gráficos e as páginas da Lara são imagem."""
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

import os
ROOT = Path(__file__).resolve().parents[2]
SP = Path(os.environ.get("DECK_WORK", Path(__file__).resolve().parent / "_work"))
SP.mkdir(parents=True, exist_ok=True)
FIG = SP / "figs_exec_v2"
F = "Arial"
INK = RGBColor(0x22, 0x30, 0x3C); GREEN = RGBColor(0x2E, 0x9E, 0x6E)
AMBER = RGBColor(0xE0, 0xA2, 0x1E); RED = RGBColor(0xCC, 0x3B, 0x3B)
LINE = RGBColor(0xE2, 0xE6, 0xEA); SOFT = RGBColor(0xF5, 0xF7, 0xF9)
GREY = RGBColor(0x5B, 0x73, 0x84); FAINT = RGBColor(0x9A, 0xA7, 0xB0); WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]; SW, SH = prs.slide_width, prs.slide_height


def slide():
    return prs.slides.add_slide(BLANK)


def _f(run, size, bold, color, italic=False):
    run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color; run.font.name = F


def txt(s, L, T, W, H, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True, spacing=1.0):
    tb = s.shapes.add_textbox(L, T, W, H); tf = tb.text_frame
    tf.word_wrap = wrap; tf.vertical_anchor = anchor
    for i, ln in enumerate(lines):
        text, size, bold, color = ln[0], ln[1], ln[2], ln[3]
        sa = ln[4] if len(ln) > 4 else 4
        ital = ln[5] if len(ln) > 5 else False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sa); p.space_before = Pt(0); p.line_spacing = spacing
        r = p.add_run(); r.text = text; _f(r, size, bold, color, ital)
    return tb


def _strip_style(shp):
    """Remove a referência de estilo do tema (sombra/efeito/preenchimento herdados) para que
    qualquer visualizador (PowerPoint, Canva, Google Slides) use só o spPr explícito — sem sombra."""
    el = shp._element
    st = el.find(qn("p:style"))
    if st is not None:
        el.remove(st)


def box(s, L, T, W, H, fill, rounded=False, line=None, lw=1.0):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, L, T, W, H)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    _strip_style(shp)
    return shp


def fit(s, path, L, T, W, H):
    iw, ih = Image.open(path).size; ar = iw / ih; bar = W / H
    if ar > bar:
        w = W; h = int(W / ar)
    else:
        h = H; w = int(H * ar)
    s.shapes.add_picture(str(path), L + (W - w) // 2, T + (H - h) // 2, w, h)


def chip(s, L, T, label, color, w=Inches(1.8)):
    c = box(s, L, T, w, Inches(0.34), color, rounded=True)
    tf = c.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label; _f(r, 11, True, WHITE)
    return c


def callout(s, L, T, W, H, text, color=GREEN):
    box(s, L, T, W, H, SOFT, rounded=False)
    box(s, L, T, Inches(0.07), H, color)
    txt(s, L + Inches(0.25), T, W - Inches(0.4), H, [(text, 14, True, INK, 0)], anchor=MSO_ANCHOR.MIDDLE)


def grid_table(s, left, top, width, col_ws, header, body, row_h,
               header_fill=SOFT, header_txt=INK, cell_styler=None):
    """'Tabela' desenhada com textboxes + linhas finas — robusta em qualquer visualizador
    (não usa a tabela nativa do OOXML, que o Canva/Google Slides remontam torta)."""
    widths = [Emu(int(width * w)) for w in col_ws]
    xs = []; x = left
    for w in widths:
        xs.append(x); x = Emu(x + w)
    y = top
    if header:
        box(s, left, y, width, row_h, header_fill)
        for j, h in enumerate(header):
            txt(s, Emu(xs[j] + Inches(0.14)), y, Emu(widths[j] - Inches(0.24)), row_h,
                [(h, 11.5, True, header_txt, 0)], anchor=MSO_ANCHOR.MIDDLE)
        y = Emu(y + row_h)
    for i, row in enumerate(body):
        for j, v in enumerate(row):
            size, bold, color = cell_styler(i, j, v) if cell_styler else (11.5, j == 0, INK)
            txt(s, Emu(xs[j] + Inches(0.14)), y, Emu(widths[j] - Inches(0.24)), row_h,
                [(v, size, bold, color, 0)], anchor=MSO_ANCHOR.MIDDLE)
        if i < len(body) - 1:
            box(s, left, Emu(y + row_h - Pt(0.4)), width, Pt(0.8), LINE)
        y = Emu(y + row_h)
    return y


def frame(kicker, title, page):
    s = slide()
    txt(s, Inches(0.55), Inches(0.42), Inches(11), Inches(0.4), [(kicker.upper(), 11, True, GREEN, 0)])
    txt(s, Inches(0.5), Inches(0.74), Inches(12.3), Inches(0.7), [(title, 29, True, INK, 0)])
    box(s, Inches(0.55), Inches(1.5), Inches(0.75), Pt(3.2), GREEN)
    box(s, Inches(0.55), Inches(7.02), Inches(12.23), Pt(0.8), LINE)
    txt(s, Inches(0.55), Inches(7.08), Inches(8), Inches(0.3),
        [("SimPred · Manutenção Preditiva · Transpetro", 9, False, FAINT, 0)])
    txt(s, Inches(11.8), Inches(7.08), Inches(1), Inches(0.3),
        [(page, 9, False, FAINT, 0)], align=PP_ALIGN.RIGHT)
    return s


def img_panel(kicker, title, img, page, chip_lbl, chip_col, callout_txt, bullets, callout_col=GREEN):
    s = frame(kicker, title, page)
    px = Inches(7.35); pw = SW - px - Inches(0.55)
    fit(s, img, Inches(0.55), Inches(1.95), px - Inches(0.9), SH - Inches(2.6))
    chip(s, px, Inches(2.05), chip_lbl, chip_col, w=Inches(2.95))
    callout(s, px, Inches(2.6), pw, Inches(1.05), callout_txt, callout_col)
    lines = []
    for b in bullets:
        lines.append(("•  " + b, 14, False, INK, 10))
    txt(s, px, Inches(4.0), pw, Inches(2.6), lines, spacing=1.05)
    return s


# ───────── 1. CAPA ─────────
s = slide()
box(s, 0, 0, SW, SH, WHITE)
txt(s, Inches(0.7), Inches(1.2), Inches(10), Inches(0.4),
    [("SIMPRED · TRANSPETRO · JUNHO 2026", 12, True, GREEN, 0)])
txt(s, Inches(0.65), Inches(2.5), Inches(11), Inches(1.6), [("Manutenção Preditiva", 60, True, INK, 0)])
txt(s, Inches(0.7), Inches(4.1), Inches(10.5), Inches(1.2),
    [("Detecção de falhas em bombas a partir dos dados de operação", 20, False, GREY, 6),
     ("Francisco Colatino de Lima · Lara Fernanda Amorim A. Cavalcante", 14, False, FAINT, 0)])
# logo do projeto (lockup claro), não um "BR" ambíguo
mk = box(s, Inches(9.9), Inches(0.95), Inches(0.75), Inches(0.75), GREEN, rounded=True)
p = mk.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "◆"; _f(r, 24, True, WHITE)
txt(s, Inches(10.78), Inches(0.95), Inches(2.4), Inches(0.8),
    [("SimPred", 22, True, INK, 2), ("MANUTENÇÃO PREDITIVA", 9, True, FAINT, 0)],
    anchor=MSO_ANCHOR.MIDDLE)

# ───────── 2. PIPELINE ─────────
s = frame("Como funciona", "Do dado bruto à detecção, em 6 etapas", "02")
STEPS = [
    ("Dados Brutos", ["CSV por sensor", "amostra de 30s a 1h"], "data"),
    ("Pré-processamento", ["Remove ruído e estado parado", "resample · clip"], "filter"),
    ("Seleção de variáveis", ["Sensores relevantes por equipamento", "+ resíduo Temp~Corrente (B-4064A)"], "select"),
    ("Modelagem", ["Autoencoder (VAE) via AutoML", "treina só em operação normal"], "model"),
    ("Detecção", ["Erro > limiar calibrado por FP", "validada em falhas históricas"], "detect"),
    ("Deploy", ["Bundle pronto p/ integração", "alarme online a validar"], "deploy"),
]
n = len(STEPS); gap = Inches(0.18); cw = (SW - Inches(1.1) - gap * (n - 1)) / n
cy = Inches(2.1); ch = Inches(2.75); x = Inches(0.55)
for i, (t, bs, icon) in enumerate(STEPS):
    box(s, x, cy, cw, ch, WHITE, rounded=True, line=LINE, lw=1.0)
    box(s, x + Inches(0.22), cy + Inches(0.22), Inches(0.58), Inches(0.58), GREEN, rounded=True)
    fit(s, FIG / f"icon_{icon}.png", x + Inches(0.32), cy + Inches(0.32), Inches(0.38), Inches(0.38))
    txt(s, x + Inches(0.18), cy + Inches(0.95), cw - Inches(0.3), Inches(0.7),
        [(t, 12.5, True, INK, 0)], spacing=1.0)
    txt(s, x + Inches(0.18), cy + Inches(1.6), cw - Inches(0.3), Inches(1.1),
        [("• " + b, 9.5, False, GREY, 6) for b in bs], spacing=1.0)
    if i < n - 1:
        ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + cw + Inches(0.01),
                                cy + ch / 2 - Inches(0.12), gap - Inches(0.02), Inches(0.24))
        ar.fill.solid(); ar.fill.fore_color.rgb = GREEN; ar.line.fill.background(); ar.shadow.inherit = False; _strip_style(ar)
    x = x + cw + gap
# tabela lateral (o que varia)
ty = Inches(5.25)
txt(s, Inches(0.55), ty, Inches(6), Inches(0.35), [("O que muda por equipamento", 13, True, INK, 0)])
grid_table(s, Inches(0.55), ty + Inches(0.4), Inches(7.6), (0.21, 0.19, 0.28, 0.32),
           ("Equip.", "resample", "filtro de estado", "passo extra"),
           [("B-8802B", "5 min", "P. Descarga > 35", "—"),
            ("B-6511502A", "5 min", "Corrente > 60", "—"),
            ("B-4064A", "1 h", "Corrente > 5", "resíduo Temp~Corrente")],
           Inches(0.33), header_fill=SOFT, header_txt=INK,
           cell_styler=lambda i, j, v: (10.5, j == 0, INK))
txt(s, Inches(8.5), ty + Inches(0.4), Inches(4.3), Inches(1.6),
    [("*Mesmo esqueleto para todos; só os parâmetros mudam por equipamento. O B-4064A troca a "
      "temperatura pelo resíduo Temp~Corrente (regime térmico).", 11, False, FAINT, 0, True)],
    anchor=MSO_ANCHOR.TOP)

# ───────── 3. RESULTADOS (2 charts) ─────────
s = frame("Resultados", "2 de 8 equipamentos prontos para produção", "03")
callout(s, Inches(0.55), Inches(1.95), Inches(12.23), Inches(0.85),
        "Detectamos as 2 falhas conhecidas com dias de antecedência — e quase zero alarme falso.", GREEN)
fit(s, FIG / "res_8802.png", Inches(0.55), Inches(3.05), Inches(5.9), Inches(3.2))
fit(s, FIG / "res_6511.png", Inches(6.9), Inches(3.05), Inches(5.9), Inches(3.2))
txt(s, Inches(0.55), Inches(6.3), Inches(5.9), Inches(0.4),
    [("B-8802B · alarme 04/07/2022 · FP 0,06%", 12, True, INK, 0)], align=PP_ALIGN.CENTER)
txt(s, Inches(6.9), Inches(6.3), Inches(5.9), Inches(0.4),
    [("B-6511502A · alarme 10/05/2023 · FP 0,05%", 12, True, INK, 0)], align=PP_ALIGN.CENTER)

# ───────── 4-5. NOVOS: B-3403C e B-90001A (análise Lara, dados interpolados) ─────────
img_panel("Detecção validada · análise Lara", "B-3403C: sinal forte ~3 semanas antes da restrição",
          FIG / "trade_3403c.png", "04", "B-3403C · Detecção validada", GREEN,
          "Baseline quieto por meses e rampa clara ~3 semanas antes do evento (12/Set/2023).",
          ["Detecção pré-falha ~100% mesmo apertando o limiar (p99 → 100% @ 1% de FP).",
           "AutoML (autoencoder DENSE, 611 trials) sobre dados interpolados.",
           "Validado em 1 evento · FP ainda in-sample → validar fora da amostra e empacotar."], GREEN)
img_panel("Requer validação · análise Lara", "B-90001A: detecção defensável, mas margem estreita",
          FIG / "trade_90001a.png", "05", "B-90001A · Requer validação", AMBER,
          "Detecta ~1 mês antes da falha (28/Ago/2021), mas a separação é estreita.",
          ["No mesmo 1% de FP, detecta só 48% — o B-3403C detecta 100%.",
           "Com debounce agressivo chega a ~85% @ 1,6% de FP (melhor equilíbrio).",
           "AutoML Isolation Forest · 1 evento · validar fora da amostra."], AMBER)

# ───────── 6. QUALIDADE DE DADOS (tabela status) ─────────
s = frame("Qualidade de dados", "O que separa os bons dos fracos", "06")
callout(s, Inches(0.55), Inches(1.95), Inches(12.23), Inches(0.85),
        "A qualidade do dado de operação — não o modelo — é o que limita os demais equipamentos.", RED)
DQ = [("B-0302C", "75% dos sensores sem sinal · 93% parada", "Bloqueado"),
      ("B-4703.24001B", "76% do tempo parada", "Bloqueado"),
      ("B-4064A", "Reconstruída em 2024/25 · mancal +24°C", "Validar"),
      ("B-24001B", "Dados intermitentes · threshold instável", "Validar")]
DQ_COL = [RED, RED, AMBER, AMBER]
grid_table(s, Inches(0.55), Inches(3.15), Inches(12.23), (0.2, 0.62, 0.18),
           ("Equipamento", "Situação do dado", "Status"), DQ, Inches(0.62),
           header_fill=INK, header_txt=WHITE,
           cell_styler=lambda i, j, v: (12, True, DQ_COL[i]) if j == 2 else (12, j == 0, INK))

# ───────── 5. B-0302C ─────────
img_panel("Bloqueado por dados", "75% dos sensores não possuem sinal útil", FIG / "donut_0302.png", "07",
          "B-0302C · Bloqueado", RED, "Sem instrumentação válida, não há o que o modelo aprenda.",
          ["Só 7 de 28 canais têm sinal real.",
           "A bomba aparece parada em 93% das leituras.",
           "Pergunta à operação: sensores descomissionados ou falha de coleta?"], RED)

# ───────── 6. B-4064A ─────────
img_panel("Requer validação", "A reconstrução da bomba mudou o regime térmico", FIG / "bars_4064.png", "08",
          "B-4064A · Validar", AMBER, "O modelo de 2024 acusa falso positivo no equipamento reconstruído.",
          ["Falha 08/2024 → reconstruída na Sulzer (peças de outra bomba).",
           "Mancal passa a operar +24°C acima do regime anterior.",
           "Ação: re-baselinar a partir de 13/01/2025."], AMBER)

# ───────── 7. B-24001B (minha versão) ─────────
img_panel("Requer validação", "Sem dado estável, o alarme varia de 36 a 8.700", FIG / "bars_24001.png", "09",
          "B-24001B · Validar", AMBER, "Dados intermitentes tornam o threshold instável (análise Lara).",
          ["Falha 06/01/2025 — vibração no mancal LNA.",
           "O nº de alarmes muda ~240× só trocando o limiar.",
           "Ação: validar a densidade/continuidade da coleta."], AMBER)

# ───────── 8. DASHBOARD ─────────
s = frame("Onde estamos", "Status executivo dos 8 equipamentos", "10")
cols = [("Prontos para produção", GREEN, [("B-8802B", "Alarme dias antes · FP 0,06%"),
                                          ("B-6511502A", "Alarme dias antes · FP 0,05%")]),
        ("Requer validação", AMBER, [("B-3403C", "Sinal forte · 1 evento · a empacotar"),
                                     ("B-90001A", "Detecção defensável · margem estreita"),
                                     ("B-4064A", "Re-baselinar pós-reconstrução"),
                                     ("B-24001B", "Threshold instável · validar coleta")]),
        ("Bloqueados por dados", RED, [("B-0302C", "75% dos sensores sem sinal"),
                                       ("B-4703.24001B", "76% do tempo parada")])]
cw = Inches(3.95); gap = Inches(0.22); x = Inches(0.55); cy = Inches(2.1); ch = Inches(4.6)
for title, col, items in cols:
    box(s, x, cy, cw, ch, WHITE, rounded=True, line=LINE, lw=1.0)
    box(s, x + Inches(0.25), cy + Inches(0.25), cw - Inches(0.5), Inches(0.55), col, rounded=True)
    hp = s.shapes[-1].text_frame.paragraphs[0]; hp.alignment = PP_ALIGN.CENTER
    hr = hp.add_run(); hr.text = title; _f(hr, 13, True, WHITE)
    pitch = int(min(Inches(1.2), (ch - Inches(1.25)) / len(items)))
    iy = cy + Inches(1.05)
    for eq, ds in items:
        txt(s, x + Inches(0.3), iy, cw - Inches(0.6), Inches(0.4), [(eq, 14, True, INK, 0)])
        txt(s, x + Inches(0.3), iy + Inches(0.34), cw - Inches(0.6), Inches(0.55),
            [(ds, 11, False, GREY, 0)], spacing=1.0)
        box(s, x + Inches(0.3), iy + pitch - Inches(0.16), cw - Inches(0.6), Pt(0.8), LINE)
        iy = iy + pitch
    x = x + cw + gap

# ───────── 9. DIVISÓRIA APÊNDICE (divisor de seção: painel verde + agenda) ─────────
s = slide(); box(s, 0, 0, SW, SH, WHITE)
PALE = RGBColor(0xD8, 0xEE, 0xE4)
box(s, 0, 0, Inches(5.5), SH, GREEN)
txt(s, Inches(0.6), Inches(2.35), Inches(4.6), Inches(0.4), [("APÊNDICE", 13, True, PALE, 0)])
txt(s, Inches(0.55), Inches(2.8), Inches(4.8), Inches(1.1), [("B-24001B", 48, True, WHITE, 0)])
txt(s, Inches(0.6), Inches(4.0), Inches(4.6), Inches(0.5), [("Análise original — Lara", 18, True, WHITE, 0)])
txt(s, Inches(0.6), Inches(6.4), Inches(4.6), Inches(0.4), [("◆ SimPred", 16, True, PALE, 0)])
txt(s, Inches(6.0), Inches(1.5), Inches(6.9), Inches(0.5),
    [("Estratégias de threshold testadas (AutoML)", 18, True, INK, 0)])
APP = [("1", "Otsu por variável", "corte automático por sensor"),
       ("2", "p90 dinâmico", "8.742 alarmes — muito sensível"),
       ("3", "p99.9", "36 alarmes — conservador"),
       ("4", "p99.5 fixo", "7.892 alarmes — poucas detecções")]
iy = Inches(2.4)
for num, t, d in APP:
    cc = box(s, Inches(6.0), iy, Inches(0.5), Inches(0.5), GREEN, rounded=False)
    cp = cc.text_frame.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run(); cr.text = num; _f(cr, 16, True, WHITE)
    txt(s, Inches(6.7), iy - Inches(0.05), Inches(6.2), Inches(0.7),
        [(t, 16, True, INK, 2), (d, 12.5, False, GREY, 0)], spacing=1.0)
    iy = iy + Inches(0.82)
callout(s, Inches(6.0), Inches(5.95), Inches(6.9), Inches(0.95),
        "O nº de alarmes varia ~240× conforme o limiar — por isso a recomendação de validar os dados.", AMBER)

# ───────── 10-13. ESTRATÉGIAS DE THRESHOLD (B-24001B, Lara) no nosso template ─────────
LARA = [
    (2, "Estratégia 1 — Otsu por variável",
     "Threshold automático (Otsu) em cada uma das 6 variáveis de vibração.",
     ["Define um corte por sensor, sem ajuste manual.", "Ponto de partida do estudo."]),
    (3, "Estratégia 2 — p90 dinâmico: 8.742 alarmes",
     "Muito sensível: detecta antes do desligamento, mas dispara demais.",
     ["Debounce 1/1 (sem suavização).", "Excesso de alarmes ao longo da série."]),
    (4, "Estratégia 3 — p99.9: 36 alarmes",
     "Conservador: detecção reduzida, comportamento parecido.",
     ["Refinamento dos parâmetros.", "Poucos alarmes — risco de perder sinal."]),
    (5, "Estratégia 4 — p99.5 fixo: 7.892 alarmes",
     "Threshold fixo (debounce 10/12): poucas detecções próximas à falha.",
     ["Resultado semelhante ao dinâmico.", "Confirma a instabilidade conforme o limiar."]),
]
for k, (p, title, call, bs) in enumerate(LARA):
    img_panel("Apêndice · B-24001B (Lara)", title, SP / f"lara_chart_{p}.png", f"{11 + k:02d}",
              "B-24001B · Validar", AMBER, call, bs, AMBER)

out = ROOT / "deploy" / "slides_executivo_transpetro_editavel_v2.pptx"
prs.save(out)
print("escrito:", out, f"({out.stat().st_size//1024} KB) · {len(prs.slides._sldIdLst)} slides")
