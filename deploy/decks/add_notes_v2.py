"""
Adiciona NOTAS DO APRESENTADOR ao deck executivo v2 e gera uma folha de cola HTML.
- Embute as notas no painel de notas de cada slide (modo apresentador do PowerPoint).
- Gera deploy/notas_apresentador_v2.html (vira PDF via Chrome headless).

Reexecutável: roda em cima do pptx editável v2; reescreve as notas a cada execução.
"""
from pathlib import Path
from html import escape

from pptx import Presentation

DEPLOY = Path(__file__).resolve().parents[2] / "deploy"
PPTX = DEPLOY / "slides_executivo_transpetro_editavel_v2.pptx"
HANDOUT = DEPLOY / "notas_apresentador_v2.html"

# (titulo curto do slide, mensagem central, [pontos de fala], tempo, [se perguntarem])
NOTES = [
    (
        "Capa",
        "Abrir com confiança e situar o projeto em uma frase.",
        [
            "\"Bom dia. Vou apresentar o SimPred, nosso trabalho de manutenção preditiva para as bombas da Transpetro.\"",
            "\"A ideia é simples: usar os próprios dados de operação que vocês já coletam para detectar falhas com dias de antecedência — sem instalar sensor novo.\"",
            "Apresente você e a Lara.",
        ],
        "~30s",
        [],
    ),
    (
        "Como funciona — pipeline de 6 etapas",
        "Temos um processo de ponta a ponta, repetível; o que muda por equipamento são só parâmetros.",
        [
            "Percorra as 6 etapas da esquerda para a direita: Dados Brutos (CSV por sensor) → Pré-processamento (tira ruído, descarta bomba parada) → Seleção de variáveis (sensores relevantes) → Modelagem (autoencoder VAE, aprende só a operação normal) → Detecção (erro acima do limiar calibrado por falso positivo) → Deploy (bundle pronto p/ integrar).",
            "Aponte a tabela: \"o esqueleto é o mesmo para todos; muda só o resample, o filtro de estado e, no B-4064A, um passo extra de resíduo Temp~Corrente.\"",
        ],
        "~2min",
        [
            "\"Vocês fazem feature engineering?\" → Seja honesto: \"não criamos features complexas. Selecionamos os sensores relevantes e, só no B-4064A, somamos um resíduo Temp~Corrente. O modelo aprende o resto sozinho.\"",
        ],
    ),
    (
        "Resultados — 2 de 8 prontos (SLIDE DE VENDA)",
        "Nos equipamentos com dado bom, o modelo detectou a falha conhecida com dias de antecedência e quase zero alarme falso.",
        [
            "Aponte o \"1º alarme\" e a linha \"falha\" à direita. \"No B-8802B o primeiro alarme dispara em 04/07, a falha é em 06/07 — dias de antecedência, com FP de 0,06%.\"",
            "\"No B-6511502A, mesma história: alarme em 10/05, FP de 0,05%.\"",
            "Explique o eixo: \"é o erro de reconstrução. Quando sobe e cruza o limiar tracejado, é o modelo dizendo: isto não parece operação normal.\"",
            "Frase de efeito: \"validamos a detecção nas falhas que já aconteceram.\"",
        ],
        "~2min",
        [
            "\"Quantos dias de antecedência?\" → No B-6511502A subimos o limiar para zerar falso positivo, o que dá ~4–5 dias de antecedência. Trade-off consciente: menos antecedência, zero alarme falso.",
        ],
    ),
    (
        "B-3403C — detecção validada (análise Lara)",
        "Mais um equipamento com sinal forte — reforça que, com dado bom, o modelo detecta.",
        [
            "\"Este é o B-3403C, analisado pela Lara. O modelo ficou quieto por meses e disparou uma rampa clara ~3 semanas antes do evento de restrição, em 12/Set/2023.\"",
            "Aponte o gráfico: baseline plano e limpo, depois um cluster denso de alertas subindo até o evento.",
            "Número honesto: \"no trade-off, a p99 pega 100% da rampa com 1% de falso positivo.\"",
            "Seja transparente: \"é 1 evento, em dados interpolados, e o FP ainda é medido in-sample — por isso está como 'detecção validada, a empacotar', não 'pronto'.\"",
        ],
        "~1,5min",
        [
            "\"É melhor que os dois prontos?\" → O sinal é tão limpo quanto, mas falta o bundle empacotado e validação fora da amostra. Por isso não subi para 'pronto'.",
        ],
    ),
    (
        "B-90001A — detecção defensável, margem estreita",
        "O sinal existe, mas é ruidoso — honestidade sobre o limite.",
        [
            "\"B-90001A teve uma falha real — afrouxamento de parafusos do mancal, em 28/Ago/2021.\"",
            "\"O modelo detecta ~1 mês antes, mas o baseline é ruidoso: o score oscila e passa perto do limiar o tempo todo.\"",
            "Número honesto: \"no melhor equilíbrio, ~85% de detecção a 1,6% de falso positivo. Apertando o limiar, a detecção cai (p99 → 48%).\"",
            "IMPORTANTE — não cite '98%': o número do cabeçalho do AutoML não reproduz nas células; use o 85% @ 1,6%.",
        ],
        "~1,5min",
        [
            "\"Dá pra colocar em produção?\" → Precisa de validação fora da amostra; a margem de separação é estreita e temos só 1 evento.",
        ],
    ),
    (
        "Qualidade de dados — bons vs fracos",
        "O que separa sucesso de fracasso é a qualidade do dado, não o modelo.",
        [
            "Leia a tabela rapidamente. \"Onde o dado é bom, funciona. Onde não é, fica bloqueado ou precisa de validação.\"",
            "\"B-0302C e B-4703 estão bloqueados por dado; B-4064A e B-24001B precisam de validação.\"",
            "Transição: \"vou detalhar cada caso nos próximos slides.\"",
        ],
        "~1min",
        [],
    ),
    (
        "B-0302C — 75% dos sensores sem sinal (BLOQUEADO)",
        "Sem instrumentação válida, não há o que o modelo aprenda.",
        [
            "\"No B-0302C, 75% dos sensores não têm sinal útil — só 7 de 28 canais. E a bomba aparece parada em 93% das leituras.\"",
            "\"Não é o modelo que falha; falta dado.\"",
            "Pergunta à operação (aponte o bullet): \"esses sensores foram descomissionados ou é falha de coleta?\" — é exatamente o que queremos confirmar com vocês.",
        ],
        "~1min",
        [],
    ),
    (
        "B-4064A — mudança de regime térmico (VALIDAR)",
        "O \"falso positivo\" do B-4064A é real: a bomba mudou fisicamente.",
        [
            "\"O modelo treinado em 2024 começou a acusar anomalia em 2025/2026. Fomos à base de falhas: a bomba falhou em 08/2024 e foi reconstruída na Sulzer, com peças de outra bomba.\"",
            "Aponte as barras: \"o mancal passou a operar +24°C acima — de 47°C para 71°C. O modelo está certo: o equipamento não é mais o mesmo.\"",
            "Ação: \"basta re-baselinar a partir de 13/01/2025 — reensinar o modelo no novo normal.\"",
            "Ponto forte: \"isso prova que o modelo é sensível a mudança real de condição.\"",
        ],
        "~1,5min",
        [],
    ),
    (
        "B-24001B — alarme varia de 36 a 8.700 (VALIDAR)",
        "Com dado intermitente, o número de alarmes é instável — não dá para confiar no limiar.",
        [
            "\"Este foi analisado pela Lara. Houve falha real em 06/01/2025, vibração no mancal LNA.\"",
            "Aponte as barras: \"mas a coleta é intermitente, e isso deixa o threshold instável: só mudando o limiar, o nº de alarmes vai de 36 a quase 8.700 — 240×.\"",
            "Ação: \"antes de produção, validar a densidade e a continuidade da coleta.\"",
        ],
        "~1,5min",
        [],
    ),
    (
        "Status executivo dos 6 equipamentos (DASHBOARD)",
        "O placar do projeto em uma tela — slide de decisão.",
        [
            "\"Resumindo: dois prontos para produção — B-8802B e B-6511502A. Dois que precisam de validação — B-4064A (re-baselinar) e B-24001B (validar coleta). E dois bloqueados por dado — B-0302C e B-4703.\"",
            "Pergunta de fechamento: \"onde vocês querem que a gente avance primeiro?\"",
            "Este é o fim natural da parte principal — o que vem depois é apêndice opcional.",
        ],
        "~1min",
        [],
    ),
    (
        "Apêndice — divisor (B-24001B / análise da Lara)",
        "Transição: a partir daqui é material detalhado e opcional.",
        [
            "\"Para quem quiser aprofundar, a Lara testou quatro estratégias de threshold no B-24001B via AutoML. Passo rápido, mas posso detalhar qualquer uma.\"",
            "Cite as 4: Otsu por variável, p90 dinâmico (8.742 alarmes), p99.9 (36), p99.5 fixo (7.892).",
            "Frase-chave: \"o nº de alarmes varia ~240× conforme o limiar — por isso recomendamos validar os dados.\"",
        ],
        "~45s",
        [],
    ),
    (
        "Apêndice 1 — Otsu por variável",
        "Ponto de partida: corte automático por sensor.",
        [
            "\"Otsu define um corte automático por sensor, sem ajuste manual, nas 6 variáveis de vibração. Foi o ponto de partida do estudo.\"",
        ],
        "~30s",
        [],
    ),
    (
        "Apêndice 2 — p90 dinâmico (8.742 alarmes)",
        "Muito sensível: detecta cedo, mas dispara demais.",
        [
            "\"p90 dinâmico é muito sensível — detecta antes do desligamento, mas dispara demais: 8.742 alarmes, sem suavização (debounce 1/1). Inviável na operação.\"",
        ],
        "~30s",
        [],
    ),
    (
        "Apêndice 3 — p99.9 (36 alarmes)",
        "O oposto: conservador, com risco de perder sinal.",
        [
            "\"p99.9 é o oposto: conservador, só 36 alarmes. Detecção reduzida, risco de perder o sinal da falha.\"",
        ],
        "~30s",
        [],
    ),
    (
        "Apêndice 4 — p99.5 fixo (7.892 alarmes)",
        "Confirma a instabilidade: o resultado depende demais do limiar.",
        [
            "\"p99.5 fixo, com debounce 10/12, dá resultado parecido com o dinâmico — confirma que, com dado instável, tudo depende do limiar. Daí a recomendação de validar a coleta antes de produção.\"",
            "FECHAMENTO: \"recapitulando o placar — 2 prontos, 2 a validar, 2 bloqueados. Fico à disposição para perguntas.\"",
        ],
        "~45s",
        [],
    ),
]


def embed_notes():
    prs = Presentation(str(PPTX))
    slides = list(prs.slides)
    assert len(slides) == len(NOTES), f"{len(slides)} slides vs {len(NOTES)} notas"
    for slide, (title, msg, talk, t, qa) in zip(slides, NOTES):
        parts = [f"MENSAGEM CENTRAL: {msg}", "", "O QUE FALAR:"]
        parts += [f"  • {b}" for b in talk]
        if qa:
            parts += ["", "SE PERGUNTAREM:"]
            parts += [f"  • {b}" for b in qa]
        parts += ["", f"Tempo sugerido: {t}"]
        slide.notes_slide.notes_text_frame.text = "\n".join(parts)
    prs.save(str(PPTX))
    print(f"✓ notas embutidas em {PPTX.name} ({len(slides)} slides)")


def build_handout():
    css = """
    * { box-sizing: border-box; }
    body { font-family: Arial, Helvetica, sans-serif; color:#1f2a30; margin:0; }
    .page { padding: 26px 34px; }
    h1 { font-size: 26px; margin:0 0 2px; color:#16past; }
    .sub { color:#6b7a82; font-size:13px; margin-bottom:18px; }
    .card { border:1px solid #e3e8ea; border-radius:10px; padding:14px 18px; margin:0 0 12px;
            page-break-inside:avoid; }
    .head { display:flex; align-items:baseline; gap:10px; margin-bottom:6px; }
    .num { background:#1f9d6a; color:#fff; font-weight:700; border-radius:6px;
           min-width:26px; height:26px; display:inline-flex; align-items:center;
           justify-content:center; font-size:13px; }
    .title { font-size:15px; font-weight:700; }
    .time { margin-left:auto; color:#6b7a82; font-size:12px; white-space:nowrap; }
    .msg { background:#f3f6f5; border-left:3px solid #1f9d6a; padding:7px 11px;
           border-radius:4px; font-size:12.5px; font-weight:600; margin:6px 0 8px; }
    .lbl { font-size:10.5px; letter-spacing:.06em; color:#1f9d6a; font-weight:700;
           text-transform:uppercase; margin:6px 0 2px; }
    ul { margin:2px 0 4px; padding-left:18px; }
    li { font-size:12.5px; line-height:1.45; margin:3px 0; }
    .qa li { color:#9a6a00; }
    """
    cards = []
    for i, (title, msg, talk, t, qa) in enumerate(NOTES, 1):
        talk_html = "".join(f"<li>{escape(b)}</li>" for b in talk)
        qa_html = ""
        if qa:
            qa_html = ('<div class="lbl">Se perguntarem</div><ul class="qa">'
                       + "".join(f"<li>{escape(b)}</li>" for b in qa) + "</ul>")
        cards.append(f"""
        <div class="card">
          <div class="head"><span class="num">{i:02d}</span>
            <span class="title">{escape(title)}</span>
            <span class="time">{escape(t)}</span></div>
          <div class="msg">{escape(msg)}</div>
          <div class="lbl">O que falar</div>
          <ul>{talk_html}</ul>
          {qa_html}
        </div>""")
    html = f"""<!doctype html><html lang="pt-br"><head><meta charset="utf-8">
    <style>{css.replace('#16past','#16322a')}</style></head><body><div class="page">
    <h1>Notas do apresentador — SimPred · Transpetro</h1>
    <div class="sub">Deck executivo v2 · 13 slides · junho 2026 · Francisco &amp; Lara</div>
    {''.join(cards)}
    </div></body></html>"""
    HANDOUT.write_text(html, encoding="utf-8")
    print(f"✓ folha de cola: {HANDOUT.name}")


if __name__ == "__main__":
    embed_notes()
    build_handout()
